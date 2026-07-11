"""여러 조건 세트 배치 백테스트 → 세트별 상세 슬라이드 + 비교 요약 PPT.

- 기준(baseline) 조건에서 세트별로 지정한 파라미터만 바꿔 run_baseline() 실행.
- raw CSV 는 한 번만 로드(load_frames)해 모든 세트가 재사용 (실행 시간 절약).
- PPT 구성: 표지 → 비교 요약 표 → 비교 차트 → 세트별 상세 7장(표지~주의사항).
- 항상 바탕화면에 .pptx 저장.

실행:
    cd c:\\cursor\\02_maxV2
    $env:PYTHONIOENCODING='utf-8'
    python build_backtest_ppt_sets.py
"""

from __future__ import annotations

from pathlib import Path

import matplotlib
import pandas as pd

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

import build_backtest_ppt as bp
import repair_v13_baseline as bl
from build_backtest_ppt import (
    BLUE,
    LIGHT,
    NAVY,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    RGBColor,
    _bg,
    _blank,
    _box,
    _build_caveats,
    _build_cond_rows,
    _cost_pct,
    _mcap_span,
    _sell_label,
    _title,
    _buy_base_label,
    _breakout_label,
    _instrument_label,
    _liquidity_label,
    _vol_label,
    add_detail_slides,
    make_charts,
)

plt.rcParams["font.family"] = "Malgun Gothic"
plt.rcParams["axes.unicode_minus"] = False

# --- 공통(기준) 실행 설정 ---
COMMON = dict(start="2020-01-01", end="2026-05-31", data_dir="data/raw")
COMMON_ARCHIVE = dict(
    start="2020-01-01",
    end="2026-05-31",
    archive_base=bl.ARCHIVE_BASE,
)

# --- 프리셋: 세트 목록 · PPT 파일명 · 표지 부제 ---
PRESETS: dict[str, dict] = {
    "volume": {
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준\n당일300%",
                "sub": "당일 거래대금 ≥ 5일평균 × 300% (룩어헤드 포함, 비교 기준)",
                "over": {},
            },
            {
                "key": "a1",
                "name": "A1 · 거래대금 전일 시프트",
                "short": "A1\n전일300%",
                "sub": "전일 거래대금 폭발 + 당일 돌파 (point-in-time, 룩어헤드 제거)",
                "over": {"vol_lag": 1},
            },
            {
                "key": "a2",
                "name": "A2 · 거래대금 필터 제거",
                "short": "A2\n필터제거",
                "sub": "거래대금 조건 미적용 — 시총·돌파만으로 진입 (신호 급증 예상)",
                "over": {"use_volume": False},
            },
            {
                "key": "c",
                "name": "C · 거래대금 600% 강화",
                "short": "C\n당일600%",
                "sub": "당일 거래대금 ≥ 5일평균 × 600% (신호 급감 예상)",
                "over": {"vol_mult": 6.0},
            },
        ],
        "out": "repair_v13_volume_sets_compare.pptx",
        "subtitle": "거래대금 조건 변화 (기준 · A1 전일 · A2 제거 · C 600%) · 뱅크롤 1억 복리",
        "chart_tag": "volume",
    },
    "base_adj": {
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준\n300%·1%",
                "sub": "K=0.7 · 시총 상위10% · 당일 거래대금 300% · 익일 시가 · 비용 1%",
                "over": {},
            },
            {
                "key": "vol600",
                "name": "세트 1 · 거래대금 600%",
                "short": "600%\n비용1%",
                "sub": "거래대금 5일평균 대비 600% (나머지 기준 동일, 비용 1%)",
                "over": {"vol_mult": 6.0},
            },
            {
                "key": "vol600_cost05",
                "name": "세트 2 · 600% + 비용 0.5%",
                "short": "600%\n비용0.5%",
                "sub": "거래대금 600% + 익일 시가 매도 비용 0.5% (나머지 기준 동일)",
                "over": {"vol_mult": 6.0, "cost_mult": 0.995},
            },
        ],
        "out": "repair_v13_base_adj_sets_compare.pptx",
        "subtitle": "기준 · 거래대금 600% · 600%+비용0.5% 비교 · 뱅크롤 1억 복리",
        "chart_tag": "base_adj",
    },
    "entry_mcap": {
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준",
                "sub": "K=0.7 · 시총 상위10% · 당일 시가 기준가 · MA5 · 당일 거래대금 300% · 익일 시가",
                "over": {},
            },
            {
                "key": "buy_close",
                "name": "세트 1 · 매수 기준가 당일 종가",
                "short": "기준가\n당일종가",
                "sub": "고가≥시가+(전일고−저)×K 충족 시 당일 종가 매수 (나머지 기준 동일)",
                "over": {"buy_price_mode": "close", "require_breakout": True},
            },
            {
                "key": "no_ma5",
                "name": "세트 2 · 시장필터 제외",
                "short": "MA5\n제외",
                "sub": "전일 종가 > MA5 시장 필터 미적용 (나머지 기준 동일)",
                "over": {"use_ma5": False},
            },
            {
                "key": "mcap_bottom",
                "name": "세트 3 · 시총 하위 10%",
                "short": "시총\n하위10%",
                "sub": "매매 전날 cross-section 시총 하위 10% (나머지 기준 동일)",
                "over": {"mcap_side": "bottom"},
            },
        ],
        "out": "repair_v13_entry_mcap_sets_compare.pptx",
        "subtitle": "기준 · 매수기준가(종가) · MA5제외 · 시총하위10% 비교 · 뱅크롤 1억 복리",
        "chart_tag": "entry_mcap",
    },
    "close_buy": {
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준\n돌파가",
                "sub": "돌파 조건 충족 시 돌파가(target) 매수 · 익일 시가 매도",
                "over": {},
            },
            {
                "key": "break_close",
                "name": "세트 1 · 돌파 후 종가 매수",
                "short": "돌파+\n종가매수",
                "sub": "고가≥시가+(전일고−저)×K 충족 시 당일 종가로 매수 (나머지 기준 동일)",
                "over": {"buy_price_mode": "close", "require_breakout": True},
            },
            {
                "key": "no_break_close",
                "name": "세트 2 · 돌파 없이 종가 매수",
                "short": "돌파X\n종가매수",
                "sub": "돌파 조건 없이 거래대금·시총·MA5 충족 시 당일 종가 매수",
                "over": {"buy_price_mode": "close", "require_breakout": False},
            },
        ],
        "out": "repair_v13_close_buy_sets_compare.pptx",
        "subtitle": "기준 · 돌파+종가매수 · 돌파없이 종가매수 비교 · 뱅크롤 1억 복리",
        "chart_tag": "close_buy",
    },
    "vol_sweep": {
        "sets": [
            {
                "key": "vol100",
                "name": "거래대금 100%",
                "short": "100%",
                "sub": "당일 거래대금 ≥ 5일평균 × 100% (나머지 기준 동일)",
                "over": {"vol_mult": 1.0},
            },
            {
                "key": "vol200",
                "name": "거래대금 200%",
                "short": "200%",
                "sub": "당일 거래대금 ≥ 5일평균 × 200% (나머지 기준 동일)",
                "over": {"vol_mult": 2.0},
            },
            {
                "key": "vol300",
                "name": "기준 · 거래대금 300%",
                "short": "300%\n(기준)",
                "sub": "당일 거래대금 ≥ 5일평균 × 300% — baseline (나머지 기준 동일)",
                "over": {},
            },
            {
                "key": "vol400",
                "name": "거래대금 400%",
                "short": "400%",
                "sub": "당일 거래대금 ≥ 5일평균 × 400% (나머지 기준 동일)",
                "over": {"vol_mult": 4.0},
            },
            {
                "key": "vol500",
                "name": "거래대금 500%",
                "short": "500%",
                "sub": "당일 거래대금 ≥ 5일평균 × 500% (나머지 기준 동일)",
                "over": {"vol_mult": 5.0},
            },
            {
                "key": "vol600",
                "name": "거래대금 600%",
                "short": "600%",
                "sub": "당일 거래대금 ≥ 5일평균 × 600% (나머지 기준 동일)",
                "over": {"vol_mult": 6.0},
            },
        ],
        "out": "repair_v13_vol_sweep_sets_compare.pptx",
        "subtitle": "거래대금 100% · 200% · 300%(기준) · 400% · 500% · 600% 비교 · 뱅크롤 1억 복리",
        "chart_tag": "vol_sweep",
    },
    "vol_bands": {
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준\n≥300%",
                "sub": "당일 거래대금 ≥ 5일평균 × 300% (기존 기준, 비교용)",
                "over": {},
            },
            {
                "key": "b50_100",
                "name": "거래대금 50%~100%",
                "short": "50~\n100%",
                "sub": "5일평균 × 50% ≤ 당일 거래대금 < × 100%",
                "over": {"vol_mult": 0.5, "vol_mult_hi": 1.0},
            },
            {
                "key": "b100_200",
                "name": "거래대금 100%~200%",
                "short": "100~\n200%",
                "sub": "5일평균 × 100% ≤ 당일 거래대금 < × 200%",
                "over": {"vol_mult": 1.0, "vol_mult_hi": 2.0},
            },
            {
                "key": "b200_300",
                "name": "거래대금 200%~300%",
                "short": "200~\n300%",
                "sub": "5일평균 × 200% ≤ 당일 거래대금 < × 300%",
                "over": {"vol_mult": 2.0, "vol_mult_hi": 3.0},
            },
            {
                "key": "b300_400",
                "name": "거래대금 300%~400%",
                "short": "300~\n400%",
                "sub": "5일평균 × 300% ≤ 당일 거래대금 < × 400%",
                "over": {"vol_mult": 3.0, "vol_mult_hi": 4.0},
            },
            {
                "key": "b400_500",
                "name": "거래대금 400%~500%",
                "short": "400~\n500%",
                "sub": "5일평균 × 400% ≤ 당일 거래대금 < × 500%",
                "over": {"vol_mult": 4.0, "vol_mult_hi": 5.0},
            },
            {
                "key": "b500_600",
                "name": "거래대금 500%~600%",
                "short": "500~\n600%",
                "sub": "5일평균 × 500% ≤ 당일 거래대금 < × 600%",
                "over": {"vol_mult": 5.0, "vol_mult_hi": 6.0},
            },
            {
                "key": "b600p",
                "name": "거래대금 600% 이상",
                "short": "≥600%",
                "sub": "당일 거래대금 ≥ 5일평균 × 600%",
                "over": {"vol_mult": 6.0},
            },
        ],
        "out": "repair_v13_vol_bands_sets_compare.pptx",
        "subtitle": "기준(≥300%) · 거래대금 구간별(50~100%…600%+) 비교 · 뱅크롤 1억 복리",
        "chart_tag": "vol_bands",
    },
    "breakout_range": {
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준\n300%",
                "sub": "K=0.7 · 시총 상위10% · MA5 · 당일 거래대금 300% · 익일 시가 · 비용 1%",
                "over": {},
            },
            {
                "key": "no_vol_range",
                "name": "세트2 · 거래대금 제거 + 기준가 구간",
                "short": "거래대금X\n구간필터",
                "sub": "거래대금 필터 제거 · 고가≥기준가 · 저가<기준가<고가 (나머지 기준 동일)",
                "over": {"use_volume": False, "require_target_in_range": True},
            },
        ],
        "out": "repair_v13_breakout_range_sets_compare.pptx",
        "subtitle": "기준 · 거래대금 제거+기준가 구간(저가<기준가<고가) 비교 · 뱅크롤 1억 복리",
        "chart_tag": "breakout_range",
        "range_analysis": True,
    },
    "prev_high_range": {
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준\n300%",
                "sub": "K=0.7 · 시총 상위10% · MA5 · 당일 거래대금 300% · 익일 시가 · 비용 1%",
                "over": {},
            },
            {
                "key": "prev_high",
                "name": "세트2 · 전일고가 구간 + 거래대금 제거",
                "short": "전일고가\n거래대금X",
                "sub": "거래대금 제거 · 저가<전일고가<당일고가 추가 · 고가≥기준가 (나머지 기준 동일)",
                "over": {"use_volume": False, "require_prev_high_in_range": True},
            },
        ],
        "out": "repair_v13_prev_high_range_sets_compare.pptx",
        "subtitle": "기준 · 전일고가 구간(저가<전일고<당일고)+거래대금제거 비교 · 뱅크롤 1억 복리",
        "chart_tag": "prev_high_range",
    },
    "liquidity_etf": {
        "common": COMMON_ARCHIVE,
        "sets": [
            {
                "key": "base",
                "name": "기준 (baseline)",
                "short": "기준\n300%",
                "sub": "K=0.7 · 시총 상위10% · MA5 · 당일 거래대금 300% · 익일 시가 · 비용 1%",
                "over": {},
            },
            {
                "key": "liq50",
                "name": "세트2 · 5일평균≥50억",
                "short": "50억\n추가",
                "sub": "매수 후보: 전일 5일평균 거래대금 ≥ 50억원 추가 (나머지 기준 동일)",
                "over": {"value_ma5_min": 5_000_000_000},
            },
            {
                "key": "liq50_no_etf",
                "name": "세트3 · ETF제외+50억",
                "short": "ETF제외\n50억",
                "sub": "ETF·ETN 제외 + 전일 5일평균 거래대금 ≥ 50억원 (나머지 기준 동일)",
                "over": {
                    "value_ma5_min": 5_000_000_000,
                    "instrument_filter": "exclude_etf",
                },
            },
            {
                "key": "etf_only",
                "name": "세트4 · ETF만+50억+0.5%",
                "short": "ETF만\n0.5%",
                "sub": "ETF·ETN만 + 5일평균≥50억 + 비용 0.5% (나머지 기준 동일)",
                "over": {
                    "value_ma5_min": 5_000_000_000,
                    "instrument_filter": "etf_only",
                    "cost_mult": 0.995,
                },
            },
        ],
        "out": "repair_v13_liquidity_etf_sets_compare.pptx",
        "subtitle": "기준 · 5일평균≥50억 · ETF제외+50억 · ETF만+50억+0.5% 비교 · 00_archive · 뱅크롤 1억 복리",
        "chart_tag": "liquidity_etf",
    },
    "prev_high_buy": {
        "common": COMMON_ARCHIVE,
        "sets": [
            {
                "key": "base",
                "name": "세트1 · 기준 (baseline)",
                "short": "기준\nK0.7",
                "sub": "K=0.7 돌파가 · 시총 상위10% · MA5 · 당일 거래대금 300% · 익일 시가 · 비용 1%",
                "over": {},
            },
            {
                "key": "prev_high",
                "name": "세트2 · 전일고가 매수",
                "short": "전일고가\n+거래대금",
                "sub": "매수 기준가·체결가 = 전일 고가 (고가≥전일고가) · 거래대금 300% 유지 (나머지 기준 동일)",
                "over": {"target_mode": "prev_high"},
            },
            {
                "key": "prev_high_no_vol",
                "name": "세트3 · 전일고가 + 거래대금 제거",
                "short": "전일고가\n거래대금X",
                "sub": "전일 고가 매수 · 거래대금 필터 제거 (나머지 기준 동일)",
                "over": {"target_mode": "prev_high", "use_volume": False},
            },
            {
                "key": "no_mcap",
                "name": "세트4 · 시총 상위10% 제거",
                "short": "시총\n필터X",
                "sub": "시총 상위 10% 필터 제거 (K=0.7·거래대금 300% 등 나머지 기준 동일)",
                "over": {"use_mcap": False},
            },
        ],
        "out": "repair_v13_prev_high_buy_sets_compare.pptx",
        "subtitle": "기준 · 전일고가 매수 · 전일고가+거래대금X · 시총필터X 비교 · 00_archive · 뱅크롤 1억 복리",
        "chart_tag": "prev_high_buy",
    },
}


def _fmt_money(won: float) -> str:
    """원 → 조/억 짧은 표기."""
    if abs(won) >= 1e12:
        return f"{won / 1e12:,.2f}조"
    return f"{won / 1e8:,.1f}억"


def _fmt_pct(v: float) -> str:
    return f"{v:+,.1f}%"


# ---------------------------------------------------------------- 배치 실행
def run_all_sets(sets: list[dict], verbose: bool = True) -> list[dict]:
    frames = bl.load_frames(**COMMON, verbose=verbose)
    results: list[dict] = []
    for spec in sets:
        if verbose:
            print(f"\n▶ [{spec['name']}] 백테스트 실행...")
        res = bl.run_baseline(**COMMON, frames=frames, verbose=False, **spec["over"])
        if not res:
            print(f"❌ [{spec['name']}] 결과 없음 — 건너뜀")
            continue
        b = res["bankroll"]
        fx = res["fixed"]
        if verbose:
            print(
                f"   고정 {fx['trades']:,}건·승률 {fx['win_rate']:.2f}%·순손익 {_fmt_money(fx['net_pnl'])}"
                f" | 뱅크롤 최종 {_fmt_money(b['final_bankroll'])}·CAGR {b['cagr']:+,.1f}%·MDD {b['mdd']:.2f}%"
            )
        results.append({**spec, "res": res})
    return results


# ---------------------------------------------------------------- 비교 차트
def make_compare_chart(results: list[dict], chart_tag: str = "default") -> Path:
    bp.CHART_DIR.mkdir(exist_ok=True)
    labels = [r["short"] for r in results]
    finals = [r["res"]["bankroll"]["final_bankroll"] / 1e8 for r in results]  # 억
    cagrs = [r["res"]["bankroll"]["cagr"] for r in results]  # %
    mdds = [r["res"]["bankroll"]["mdd"] for r in results]  # %
    trades = [r["res"]["fixed"]["trades"] for r in results]  # 고정 매매건수
    avg_daily = [r["res"]["fixed"]["avg_daily"] for r in results]  # 일평균 매매건수
    over_pct = [r["res"]["signal_freq"]["over_slot_pct"] for r in results]  # 10종 초과일 %

    palette = ["#9AA6C0", "#2E5BFF", "#1B9E5A", "#E0663A", "#8E44AD"]
    colors = [palette[i % len(palette)] for i in range(len(labels))]

    fig, axes = plt.subplots(2, 3, figsize=(13.2, 7.6), dpi=150)

    def _bar(ax, vals, title, fmt, *, log=False):
        bars = ax.bar(labels, vals, color=colors)
        ax.set_title(title, fontsize=12.5, weight="bold")
        ax.grid(True, axis="y", ls="--", alpha=0.3)
        if log:
            ax.set_yscale("log")
        ax.margins(y=0.20)
        ax.tick_params(axis="x", labelsize=9)
        for b, v in zip(bars, vals):
            ax.text(
                b.get_x() + b.get_width() / 2, v, fmt(v),
                ha="center", va="bottom", fontsize=9.5, weight="bold",
            )

    _bar(axes[0][0], finals, "최종 뱅크롤 (억원, 로그)", lambda v: _fmt_money(v * 1e8), log=True)
    _bar(axes[0][1], cagrs, "CAGR (년평균 수익률, %)", lambda v: f"{v:,.0f}%")
    _bar(axes[0][2], mdds, "MDD (최대 낙폭, %)", lambda v: f"{v:.2f}%")
    _bar(axes[1][0], trades, "고정 매매 건수 (빈도)", lambda v: f"{v:,.0f}")
    _bar(axes[1][1], avg_daily, "일평균 매매 건수", lambda v: f"{v:.1f}")
    _bar(axes[1][2], over_pct, "신호 10종 초과일 비율 (%)", lambda v: f"{v:.1f}%")

    fig.suptitle("조건 세트 비교 — 성과 · 빈도 지표", fontsize=16, weight="bold")
    fig.tight_layout(rect=(0, 0, 1, 0.96))
    out = bp.CHART_DIR / f"chart_sets_compare_{chart_tag}.png"
    fig.savefig(out)
    plt.close(fig)
    return out


# ---------------------------------------------------------------- 요약 슬라이드
def add_cover(prs: Presentation, period_line: str, subtitle: str) -> None:
    s = _blank(prs)
    _bg(s, NAVY)
    bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.55), Inches(2.6), Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    _box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.5),
         "조건 세트 비교 백테스트", 40, color=WHITE, bold=True)
    _box(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.8),
         subtitle, 20, color=RGBColor(0xB9, 0xC6, 0xE6))
    _box(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
         period_line, 15, color=RGBColor(0x9A, 0xA6, 0xC0))


def add_summary_table(prs: Presentation, results: list[dict], *, common: dict) -> None:
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "세트별 비교 요약")

    names = [r["name"].replace(" · ", "\n").replace(" (baseline)", "") for r in results]
    header = ["지표"] + names

    def row(label, fn):
        return [label] + [fn(r["res"]) for r in results]

    rows = [
        header,
        row("돌파 조건", lambda x: _breakout_label(x["params"])),
        row("매수 체결가", lambda x: _buy_base_label(x["params"])),
        row("시장 필터", lambda x: "전일 종가 > MA5" if x["params"].get("use_ma5_filter", True) else "미적용"),
        row("시총 필터", lambda x: f"{_mcap_span(x['params'])} (D-1 cross-section)"),
        row("거래대금 조건", lambda x: _vol_label(x["params"])),
        row("유동성 하한", lambda x: _liquidity_label(x["params"])),
        row("종목 유형", lambda x: _instrument_label(x["params"])),
        row("거래비용", lambda x: _cost_pct(x["params"])),
        row("고정 매매건수", lambda x: f"{x['fixed']['trades']:,}건"),
        row("총 매매일수", lambda x: f"{x['fixed']['trade_days']:,}일"),
        row("일평균 매매건수", lambda x: f"{x['fixed']['avg_daily']:.2f}건"),
        row("신호 10종초과일 비율", lambda x: f"{x['signal_freq']['over_slot_pct']:.1f}%"),
        row("고정 승률", lambda x: f"{x['fixed']['win_rate']:.2f}%"),
        row("고정 순손익", lambda x: _fmt_money(x["fixed"]["net_pnl"])),
        row("뱅크롤 체결건수", lambda x: f"{x['bankroll']['trades']:,}건"),
        row("뱅크롤 승률", lambda x: f"{x['bankroll']['win_rate']:.2f}%"),
        row("최종 뱅크롤", lambda x: _fmt_money(x["bankroll"]["final_bankroll"])),
        row("누적 수익률", lambda x: f"{x['bankroll']['cum_ret']:+,.0f}%"),
        row("CAGR", lambda x: f"{x['bankroll']['cagr']:+,.1f}%"),
        row("MDD", lambda x: f"{x['bankroll']['mdd']:.2f}%"),
    ]

    n_sets = len(results)
    first_w = 2.7
    rest_w = (12.4 - first_w) / n_sets
    col_w = [Inches(first_w)] + [Inches(rest_w)] * n_sets
    tbl_font = 9.5 if n_sets > 6 else 11.5
    tbl_row_h = 0.32 if n_sets > 6 else 0.40
    bp._table(
        s, Inches(0.5), Inches(1.35), Inches(12.4), rows,
        col_w=col_w, font_size=tbl_font, row_h=tbl_row_h,
        highlight_rows={
            6: RGBColor(0xFF, 0xF3, 0xE0),   # 신호 10종초과일 비율 (빈도)
            11: RGBColor(0xE9, 0xEF, 0xFF),  # 최종 뱅크롤
            13: RGBColor(0xE6, 0xF0, 0xE9),  # CAGR
        },
    )
    eff = results[0]["res"]["params"].get("effective_end")
    _box(s, Inches(0.5), Inches(7.0), Inches(12.4), Inches(0.4),
         f"※ 기간 {common['start']} ~ {common['end']} (실제 데이터 종료 {eff}) · 각 세트는 기준 조건에서 해당 항목만 변경",
         11, color=RGBColor(0x60, 0x66, 0x72))


def add_compare_chart_slide(prs: Presentation, chart: Path) -> None:
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "세트별 성과 비교 차트")
    from PIL import Image

    with Image.open(chart) as im:
        ratio = im.width / im.height
    h = Inches(5.6)
    w = Emu(int(h * ratio))
    x = Emu(int((SLIDE_W - w) / 2))
    s.shapes.add_picture(str(chart), x, Inches(1.5), height=h)


def add_range_rejection_slide(prs: Presentation, stats: dict[str, int]) -> None:
    """세트2 구간 필터 — 저가<기준가 미충족 제외 건수."""
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "세트2 · 기준가 구간 필터 제외 건수")

    loose = stats["loose"]
    strict = stats["strict"]
    rej_low = stats["rejected_not_above_low"]
    rej_high = stats["rejected_not_below_high"]
    pct = (rej_low / loose * 100) if loose else 0.0

    rows = [
        ["구분", "건수", "설명"],
        [
            "A. loose 신호 (거래대금 제거)",
            f"{loose:,}건",
            "시총·MA5·고가≥기준가 충족 (구간 필터 없음)",
        ],
        [
            "B. strict 신호 (세트2 적용)",
            f"{strict:,}건",
            "A + 저가 < 기준가 < 고가",
        ],
        [
            "C. 저가<기준가 미충족 제외",
            f"{rej_low:,}건",
            "A 충족이나 기준가 ≤ 당일 저가 (갭상승 등)",
        ],
        [
            "C / A 비율",
            f"{pct:.2f}%",
            "전체 loose 신호 중 저가 조건으로만 제외된 비율",
        ],
        [
            "D. 기준가≥고가 제외 (참고)",
            f"{rej_high:,}건",
            "A 충족·저가<기준가이나 기준가 ≥ 고가 (경계)",
        ],
    ]
    bp._table(s, Inches(0.5), Inches(1.35), Inches(12.4), rows, col_w=[Inches(3.2), Inches(2.0), Inches(7.2)])
    _box(
        s, Inches(0.5), Inches(6.85), Inches(12.4), Inches(0.55),
        "※ 세트2는 거래대금 필터 없음 · 매수 체결은 돌파가(target) · 검증 A=loose / B=strict / C=저가 미충족",
        11, color=RGBColor(0x60, 0x66, 0x72),
    )


# ---------------------------------------------------------------- 드라이버
def build_sets_report(
    results: list[dict],
    out_name: str,
    *,
    subtitle: str,
    chart_tag: str,
    common: dict,
    range_stats: dict[str, int] | None = None,
) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    eff = results[0]["res"]["params"].get("effective_end")
    period_line = f"기간  {common['start']} ~ {common['end']}  (실제 데이터 종료: {eff})"

    print("▶ 비교 차트 생성...")
    cmp_chart = make_compare_chart(results, chart_tag=chart_tag)

    add_cover(prs, period_line, subtitle)
    add_summary_table(prs, results, common=common)
    add_compare_chart_slide(prs, cmp_chart)
    if range_stats is not None:
        add_range_rejection_slide(prs, range_stats)

    print("▶ 세트별 상세 슬라이드 생성...")
    for r in results:
        res = r["res"]
        charts = make_charts(res, tag=r["key"])
        add_detail_slides(
            prs, res, charts,
            title=r["name"],
            subtitle=r["sub"],
            cond_rows=_build_cond_rows(res["params"]),
            caveats=_build_caveats(res["params"]),
        )

    out = bp._resolve_out_path(out_name)
    prs.save(str(out))
    print(f"✅ 저장 완료: {out}")
    return out


def main() -> None:
    import argparse

    ap = argparse.ArgumentParser(description="조건 세트 배치 백테스트 → 비교 PPT")
    ap.add_argument(
        "--preset",
        choices=list(PRESETS.keys()),
        default="volume",
        help="세트 구성 프리셋 (default: volume)",
    )
    args = ap.parse_args()
    cfg = PRESETS[args.preset]
    sets = cfg["sets"]
    common = cfg.get("common", COMMON)

    print(f"▶ 조건 세트 배치 백테스트 시작 (preset={args.preset})")
    frames = bl.load_frames(**common, verbose=True)
    results: list[dict] = []
    for spec in sets:
        print(f"\n▶ [{spec['name']}] 백테스트 실행...")
        res = bl.run_baseline(**common, frames=frames, verbose=False, **spec["over"])
        if not res:
            print(f"❌ [{spec['name']}] 결과 없음 — 건너뜀")
            continue
        b = res["bankroll"]
        fx = res["fixed"]
        print(
            f"   고정 {fx['trades']:,}건·승률 {fx['win_rate']:.2f}%·순손익 {_fmt_money(fx['net_pnl'])}"
            f" | 뱅크롤 최종 {_fmt_money(b['final_bankroll'])}·CAGR {b['cagr']:+,.1f}%·MDD {b['mdd']:.2f}%"
        )
        results.append({**spec, "res": res})
    if not results:
        print("❌ 실행된 세트가 없습니다.")
        return

    range_stats = None
    if cfg.get("range_analysis"):
        start_ts = pd.Timestamp(common["start"])
        end_ts = pd.Timestamp(common["end"])
        cutoffs = bl._build_daily_cutoffs(frames, 0.1, ratio_hi=None, side="top")
        range_stats = bl.count_all_breakout_rejected_not_above_low(
            frames, cutoffs, start_ts, end_ts,
        )
        print(
            f"\n▶ 구간 필터 제외 분석: loose={range_stats['loose']:,}건 · "
            f"strict={range_stats['strict']:,}건 · "
            f"저가<기준가 미충족={range_stats['rejected_not_above_low']:,}건 "
            f"({range_stats['rejected_not_above_low'] / range_stats['loose'] * 100 if range_stats['loose'] else 0:.2f}%)"
        )

    build_sets_report(
        results,
        cfg["out"],
        subtitle=cfg["subtitle"],
        chart_tag=cfg["chart_tag"],
        common=common,
        range_stats=range_stats,
    )


if __name__ == "__main__":
    main()
