"""repair_v13 스냅샷 + 뱅크롤 복리 백테스트 결과 → PPT 리포트 생성.

- repair_v13_bankroll.run_backtest() 를 호출해 결과 dict 확보
- matplotlib 로 차트 3종 생성 (뱅크롤 성장곡선, 년도별 수익률, 고정 vs 뱅크롤)
- python-pptx 로 표지 + 조건 + 결과표 + 차트 + 주의사항 슬라이드 구성
- 바탕화면(있으면)에 .pptx 저장
"""

from __future__ import annotations

import os
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

import repair_v13_bankroll as rv

# --- 한글 폰트 ---
plt.rcParams["font.family"] = "Malgun Gothic"
plt.rcParams["axes.unicode_minus"] = False

FONT = "맑은 고딕"
NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x2E, 0x5B, 0xFF)
GREEN = RGBColor(0x1B, 0x9E, 0x5A)
RED = RGBColor(0xD1, 0x3A, 0x3A)
GRAY = RGBColor(0x60, 0x66, 0x72)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CHART_DIR = Path("_ppt_charts")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _eok(won: float) -> float:
    return won / 1e8


# ---------------------------------------------------------------- 차트
def make_charts(res: dict, tag: str = "") -> dict:
    CHART_DIR.mkdir(exist_ok=True)
    paths: dict[str, Path] = {}
    sfx = f"_{tag}" if tag else ""

    # 1) 뱅크롤 성장 곡선 (log scale, 억원)
    eq = res["equity_curve"]
    dates = [__import__("pandas").Timestamp(r["date"]) for r in eq]
    vals = [_eok(r["bankroll"]) for r in eq]
    fig, ax = plt.subplots(figsize=(10.5, 5.2), dpi=150)
    ax.plot(dates, vals, color="#2E5BFF", lw=2.2)
    ax.fill_between(dates, vals, color="#2E5BFF", alpha=0.08)
    ax.set_yscale("log")
    ax.set_title("뱅크롤 성장 곡선 (매일 아침 정산, 로그 스케일)", fontsize=15, weight="bold")
    ax.set_ylabel("뱅크롤 (억원)")
    ax.grid(True, which="both", ls="--", alpha=0.3)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    fig.autofmt_xdate()
    fig.tight_layout()
    p1 = CHART_DIR / f"chart_growth{sfx}.png"
    fig.savefig(p1)
    plt.close(fig)
    paths["growth"] = p1

    # 2) 년도별 수익률 막대
    yr = res["yearly"]
    years = [str(r["year"]) for r in yr]
    rets = [r["return_pct"] for r in yr]
    fig, ax = plt.subplots(figsize=(6.2, 5.0), dpi=150)
    palette = ["#2E5BFF", "#1B9E5A", "#D1A33A", "#8E44AD", "#E0663A", "#3AA0C0", "#C0405A"]
    colors = [palette[i % len(palette)] for i in range(len(years))]
    bars = ax.bar(years, rets, color=colors)
    ax.set_title("년도별 수익률", fontsize=15, weight="bold")
    ax.set_ylabel("수익률 (%)")
    ax.grid(True, axis="y", ls="--", alpha=0.3)
    for b, v in zip(bars, rets):
        ax.text(b.get_x() + b.get_width() / 2, v, f"+{v:,.0f}%", ha="center", va="bottom", fontsize=11, weight="bold")
    ax.margins(y=0.15)
    fig.tight_layout()
    p2 = CHART_DIR / f"chart_yearly{sfx}.png"
    fig.savefig(p2)
    plt.close(fig)
    paths["yearly"] = p2

    # 3) 고정 vs 뱅크롤 — 매매건수 & 승률
    fx, bk = res["fixed"], res["bankroll"]
    fig, (axL, axR) = plt.subplots(1, 2, figsize=(9.6, 4.6), dpi=150)
    axL.bar(["고정 1천만원", "뱅크롤 복리"], [fx["trades"], bk["trades"]], color=["#9AA6C0", "#2E5BFF"])
    axL.set_title("체결 매매 건수", fontsize=13, weight="bold")
    for i, v in enumerate([fx["trades"], bk["trades"]]):
        axL.text(i, v, f"{v:,}", ha="center", va="bottom", fontsize=11, weight="bold")
    axL.grid(True, axis="y", ls="--", alpha=0.3)
    axL.margins(y=0.15)
    axR.bar(["고정 1천만원", "뱅크롤 복리"], [fx["win_rate"], bk["win_rate"]], color=["#9AA6C0", "#1B9E5A"])
    axR.set_title("승률 (%)", fontsize=13, weight="bold")
    axR.set_ylim(0, 100)
    for i, v in enumerate([fx["win_rate"], bk["win_rate"]]):
        axR.text(i, v, f"{v:.2f}%", ha="center", va="bottom", fontsize=11, weight="bold")
    axR.grid(True, axis="y", ls="--", alpha=0.3)
    fig.suptitle("고정 1천만원 vs 뱅크롤 복리", fontsize=15, weight="bold")
    fig.tight_layout()
    p3 = CHART_DIR / f"chart_compare{sfx}.png"
    fig.savefig(p3)
    plt.close(fig)
    paths["compare"] = p3

    return paths


# ---------------------------------------------------------------- PPT helpers
def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, x, y, w, h, text, size, *, color=NAVY, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return tb


def _accent_bar(slide, y=Inches(1.15)):
    bar = slide.shapes.add_shape(1, Inches(0.6), y, Inches(2.2), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def _title(slide, text):
    _box(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.8), text, 30, color=NAVY, bold=True)
    _accent_bar(slide)


def _table(slide, x, y, w, rows: list[list[str]], *, col_w=None, header=True,
           font_size=14, highlight_rows=None, row_h=0.45):
    highlight_rows = highlight_rows or {}
    n_r, n_c = len(rows), len(rows[0])
    h = Inches(row_h * n_r)
    gtbl = slide.shapes.add_table(n_r, n_c, x, y, w, h).table
    if col_w:
        for ci, cw in enumerate(col_w):
            gtbl.columns[ci].width = cw
    for ri, row in enumerate(rows):
        gtbl.rows[ri].height = Inches(row_h)
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri, ci)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT
            r.font.size = Pt(font_size)
            if header and ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                r.font.color.rgb = WHITE
                r.font.bold = True
            elif ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_rows[ri]
                r.font.bold = True
                r.font.color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = NAVY
    return gtbl


def _pic_centered(slide, img: Path, *, top=Inches(1.5), height=Inches(5.4)):
    from PIL import Image

    with Image.open(img) as im:
        ratio = im.width / im.height
    h = height
    w = Emu(int(h * ratio))
    x = Emu(int((SLIDE_W - w) / 2))
    slide.shapes.add_picture(str(img), x, top, height=h)


# ---------------------------------------------------------------- 슬라이드
def add_detail_slides(
    prs: Presentation,
    res: dict,
    charts: dict,
    *,
    title: str,
    subtitle: str,
    cond_rows: list[list[str]],
    caveats: str,
) -> None:
    """주어진 Presentation 에 세트 상세 7장(표지~주의사항)을 추가한다."""
    p = res["params"]
    fx = res["fixed"]
    bk = res["bankroll"]

    # ---- 1) 표지
    s = _blank(prs)
    _bg(s, NAVY)
    bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.55), Inches(2.6), Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    _box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.5),
         title, 40, color=WHITE, bold=True)
    _box(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.8),
         subtitle, 22, color=RGBColor(0xB9, 0xC6, 0xE6))
    eff_note = ""
    if p.get("effective_end") and p["effective_end"] != p["end"]:
        eff_note = f"  (실제 데이터 종료: {p['effective_end']})"
    _box(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
         f"기간  {p['start']} ~ {p['end']}{eff_note}\n"
         f"전략  K={p['K']} · {_vol_label(p)} · MA5 · 시총 {_mcap_span(p)} · 매도 {_sell_label(p)} · 비용 {_cost_pct(p)}",
         15, color=RGBColor(0x9A, 0xA6, 0xC0))

    # ---- 2) 백테스트 조건
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "백테스트 조건")
    _table(s, Inches(0.9), Inches(1.5), Inches(11.5), cond_rows,
           col_w=[Inches(6.3), Inches(5.2)], font_size=14)

    # ---- 3) 고정 1천만원 결과
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "[고정 1천만원] 결과 — 매 건 1,000만원 고정")
    frows = [
        ["항목", "값"],
        ["매매 건수", f"{fx['trades']:,}건"],
        ["매매일", f"{fx['trade_days']:,}일"],
        ["순손익", f"{fx['net_pnl']:+,.0f}원"],
        ["이익", f"+{fx['profit']:,.0f}원 ({fx['profit_cnt']:,}건)"],
        ["손실", f"{fx['loss']:,.0f}원 ({fx['loss_cnt']:,}건)"],
        ["승률", f"{fx['win_rate']:.2f}%"],
        ["매매당 평균 수익률", f"+{fx['avg_ret']:.2f}%"],
        ["일평균 매수", f"{fx['avg_daily']:.2f}건/일"],
    ]
    _table(s, Inches(0.9), Inches(1.6), Inches(8.2), frows,
           col_w=[Inches(4.0), Inches(4.2)], font_size=15,
           highlight_rows={3: RGBColor(0xE6, 0xF0, 0xE9), 6: RGBColor(0xE9, 0xEF, 0xFF)})
    _box(s, Inches(9.5), Inches(1.7), Inches(3.2), Inches(4),
         "매 신호마다\n1,000만원씩\n독립 매수했다고\n가정한\n단순 합산 결과입니다.\n\n뱅크롤(자금 한도)\n제약이 없어\n실제 체결보다\n건수가 많습니다.",
         15, color=GRAY)

    # ---- 4) 뱅크롤 복리 결과
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "[뱅크롤 연동] 결과 — 시작 1억 복리")
    brows = [
        ["항목", "값"],
        ["체결 매매 건수", f"{bk['trades']:,}건"],
        ["매매일", f"{bk['trade_days']:,}일"],
        ["승률", f"{bk['win_rate']:.2f}%"],
        ["최종 뱅크롤", f"{bk['final_bankroll']:,.0f}원"],
        ["최종 누적 수익률", f"+{bk['cum_ret']:,.2f}%"],
        ["년평균 수익률 (CAGR)", f"+{bk['cagr']:,.2f}%"],
        ["MDD (최대 낙폭)", f"{bk['mdd']:.2f}%"],
    ]
    _table(s, Inches(0.9), Inches(1.6), Inches(6.3), brows,
           col_w=[Inches(3.3), Inches(3.0)], font_size=15,
           highlight_rows={4: RGBColor(0xE9, 0xEF, 0xFF), 6: RGBColor(0xE6, 0xF0, 0xE9)})

    # 매매 방식 설명 (좌측 표 아래)
    method_y = Inches(1.6) + Inches(0.45 * len(brows)) + Inches(0.15)
    _box(s, Inches(0.9), method_y, Inches(6.3), Inches(0.32),
         "▮ 매매 방식 (복리 시뮬레이션)", 14, color=BLUE, bold=True)
    if p.get("sell_mode") == "same_close":
        settle_line = "• 매일 당일 종가로 매수분 전량 매도 → 뱅크롤 재정산 (오버나이트 미보유)\n"
        exit_line = f"• 매수=돌파가, 매도=당일 종가, 매도 시 비용 {_cost_pct(p)}"
    else:
        settle_line = "• 매일 아침 전일 매수분을 시가 전량 매도 → 뱅크롤 재정산\n"
        exit_line = f"• 매수=돌파가, 매도=익일 시가, 매도 시 비용 {_cost_pct(p)}"
    method_txt = (
        f"• 시작 뱅크롤 {p['initial_bankroll']:,.0f}원, 전량 복리 재투자\n"
        + settle_line
        + f"• 재정산 뱅크롤 ÷{p['max_daily_positions']} = 종목당 매수 한도\n"
        + f"• 신호 {p['max_daily_positions']}종 초과 → 랜덤 {p['max_daily_positions']}종만 매수, 나머지 미매매 (seed={p['seed']})\n"
        + f"• 신호 {p['max_daily_positions']}종 이하 → 전부 매수 (남는 현금 이월)\n"
        + exit_line
    )
    _box(s, Inches(0.9), method_y + Inches(0.36), Inches(6.3), Inches(1.7),
         method_txt, 12, color=RGBColor(0x33, 0x3A, 0x48))

    yr = res["yearly"]
    yrows = [["연도", "매매", "승률", "수익률", "뱅크롤(종료)"]]
    for r in yr:
        yrows.append([
            str(r["year"]), f"{r['trades']:,}", f"{r['win_rate']:.1f}%",
            f"+{r['return_pct']:,.0f}%", f"{_eok(r['end_bankroll']):,.1f}억",
        ])
    _table(s, Inches(7.5), Inches(1.6), Inches(5.2), yrows,
           col_w=[Inches(0.9), Inches(1.1), Inches(1.0), Inches(1.2), Inches(1.5)],
           font_size=13)
    _box(s, Inches(7.5), Inches(1.6) + Inches(0.45 * len(yrows)) + Inches(0.15),
         Inches(5.2), Inches(0.8),
         "※ 년도별 수익률 = 해당 연 시작 뱅크롤 대비 연말 뱅크롤", 12, color=GRAY)

    # ---- 5) 차트: 뱅크롤 성장곡선
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "뱅크롤 성장 곡선")
    _pic_centered(s, charts["growth"], top=Inches(1.45), height=Inches(5.5))

    # ---- 6) 차트: 년도별 + 비교
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "년도별 수익률 · 고정 vs 뱅크롤")
    slide = s
    from PIL import Image
    with Image.open(charts["yearly"]) as im:
        r1 = im.width / im.height
    h1 = Inches(4.9)
    slide.shapes.add_picture(str(charts["yearly"]), Inches(0.7), Inches(1.6), height=h1)
    with Image.open(charts["compare"]) as im:
        r2 = im.width / im.height
    h2 = Inches(3.6)
    w2 = Emu(int(h2 * r2))
    slide.shapes.add_picture(str(charts["compare"]), Inches(6.6), Inches(2.2), height=h2)

    # ---- 7) 주의사항
    s = _blank(prs)
    _bg(s, WHITE)
    _title(s, "해석 시 주의사항")
    _box(s, Inches(0.9), Inches(1.6), Inches(11.6), Inches(5.4), caveats, 14, color=RGBColor(0x33, 0x3A, 0x48))


def build_ppt(
    res: dict,
    charts: dict,
    out_path: Path,
    *,
    title: str,
    subtitle: str,
    cond_rows: list[list[str]],
    caveats: str,
) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_detail_slides(
        prs, res, charts,
        title=title, subtitle=subtitle, cond_rows=cond_rows, caveats=caveats,
    )
    prs.save(str(out_path))


def _resolve_out_path(name: str) -> Path:
    candidates = [
        Path(os.path.expanduser("~")) / "Desktop",
        Path(os.path.expanduser("~")) / "OneDrive" / "Desktop",
        Path(os.path.expanduser("~")) / "바탕 화면",
    ]
    for d in candidates:
        if d.exists():
            return d / name
    return Path.cwd() / name


def _mcap_span(p: dict) -> str:
    ratio = p.get("mcap_top_ratio", 0.1) * 100
    hi = p.get("mcap_ratio_hi")
    if p.get("mcap_side") == "bottom":
        return f"하위 {ratio:.0f}%"
    if hi:
        return f"상위 {ratio:.0f}%~{hi*100:.0f}%"
    return f"상위 {ratio:.0f}%"


def _buy_base_label(p: dict) -> str:
    mode = p.get("buy_price_mode", "target")
    if mode == "close":
        return "당일 종가 (체결가)"
    return "돌파가 = 당일 시가 + (전일 고가 − 전일 저가) × K"


def _breakout_label(p: dict) -> str:
    if p.get("require_breakout", True):
        return "당일 고가 ≥ 시가 + (전일 고−저) × K"
    return "미적용"


def _sell_label(p: dict) -> str:
    return "당일 종가" if p.get("sell_mode") == "same_close" else "익일 시가"


def _cost_pct(p: dict) -> str:
    """거래비용(%) 표기 — 0.5% 등 소수 비용도 슬라이드 전체에서 동일하게."""
    pct = round((1 - p["cost_mult"]) * 100, 2)
    if abs(pct - round(pct)) < 1e-9:
        return f"{int(round(pct))}%"
    return f"{pct:.1f}%"


def _vol_label(p: dict) -> str:
    if not p.get("use_volume", True):
        return "미적용 (필터 제거)"
    base = "전일" if p.get("vol_lag") else "당일"
    hi = p.get("vol_mult_hi")
    lo = p.get("vol_spike_mult", 3.0)
    if hi is not None:
        return f"{base} 거래대금 5일평균 × {lo*100:.0f}%~{hi*100:.0f}%"
    return f"{base} 거래대금 ≥ 5일평균 × {lo*100:.0f}%"


def _mcap_label(p: dict) -> str:
    mode = p.get("mcap_mode", "snapshot")
    span = _mcap_span(p)
    if "point-in-time" in mode:
        return f"매매 전날 기준 일별 cross-section {span}"
    return f"스냅샷(최신 시총) {span}"


def _build_cond_rows(p: dict) -> list[list[str]]:
    period = f"{p['start']} ~ {p['end']}"
    if p.get("effective_end") and p["effective_end"] != p["end"]:
        period += f"  (실제 데이터: {p['effective_end']})"
    rows = [
        ["항목", "값"],
        ["기간", period],
        ["대상 종목 수", f"{p['symbol_count']:,}종"],
        ["시총 필터", _mcap_label(p)],
        ["돌파 조건", _breakout_label(p)],
        ["매수 체결가", _buy_base_label(p)],
        ["K값 (전일 변동폭 계수)", f"{p['K']}"],
        ["거래대금 조건", _vol_label(p)],
        ["시장 필터", "전일 종가 > 전일 MA5" if p.get("use_ma5_filter", True) else "미적용"],
        ["매도 / 거래비용", f"{_sell_label(p)} · {_cost_pct(p)}"],
        ["시작 뱅크롤", f"{p['initial_bankroll']:,.0f}원 (복리)"],
        ["슬롯 / 일 최대 종목", f"뱅크롤 ÷ {p['max_daily_positions']} · 최대 {p['max_daily_positions']}종"],
        ["신호 초과 시", f"랜덤 추출 (seed={p['seed']})"],
        ["뱅크롤 평가 시점", "매일 아침 시가 매도 후 현금"],
    ]
    return rows


def _build_caveats(p: dict) -> str:
    lines = []
    if "point-in-time" in p.get("mcap_mode", ""):
        span = _mcap_span(p)
        lines.append(
            f"1.  시총 필터 (point-in-time) — 매매 전날(D-1) 기준 일별 cross-section {span}로 필터해 "
            "'미래 시총으로 과거를 거르는' 룩어헤드를 제거했습니다. (스냅샷 방식 대비 개선된 표준 조건)"
        )
    else:
        lines.append(
            "1.  룩어헤드(look-ahead) — 시총 필터가 각 종목의 '최신(마지막 행) 시총'을 과거 전 구간에 적용해 "
            "결과가 낙관적으로 부풀 수 있습니다."
        )
    lines.append(
        "2.  생존편향 — 대상 데이터가 현재 상장·생존 종목 위주라면, 중간 상장폐지·거래정지 종목이 빠져 성과가 과대평가될 수 있습니다."
    )
    if p.get("effective_end") and p["effective_end"] != p["end"]:
        lines.append(
            f"3.  데이터 종료일 — 원본 raw 데이터가 {p['effective_end']}까지만 존재합니다. "
            f"요청 종료일({p['end']}) 이후 구간은 비어 있으며, 누적수익률·CAGR은 실제 종료일 기준으로 계산했습니다."
        )
    if not p.get("use_volume", True):
        lines.append(
            f"{len(lines)+1}.  거래대금 필터 제거 — 유동성·거래 확인 없이 시총·돌파만으로 진입해 신호수와 "
            "저유동 종목 편입이 크게 늘 수 있습니다. 실전 체결 난이도는 반영되지 않았습니다."
        )
    elif p.get("vol_lag"):
        lines.append(
            f"{len(lines)+1}.  거래대금 필터 (PIT) — 전일 거래대금 폭발을 조건으로 해 미래정보(당일 최종 거래량) 없이 "
            "판정합니다. '당일 거래대금' 방식 대비 룩어헤드가 제거된 보수적 조건입니다."
        )
    else:
        lines.append(
            f"{len(lines)+1}.  거래대금 필터 (룩어헤드) — 당일 최종 거래대금으로 장중 진입을 판정하므로, 실제 매수 시점엔 "
            "알 수 없는 미래정보가 섞여 성과가 낙관적으로 부풀 수 있습니다."
        )
    if p.get("buy_price_mode") == "close":
        lines.append(
            f"{len(lines)+1}.  매수 체결가 (당일 종가) — 돌파 여부·기타 필터 충족 시 "
            "돌파가(target)가 아닌 **당일 종가**로 매수합니다. 장 마감 후 체결 가정이며, "
            "장중 돌파 시점 체결과는 다릅니다."
        )
    if not p.get("require_breakout", True):
        lines.append(
            f"{len(lines)+1}.  돌파 조건 미적용 — 고가≥시가+(전일고−저)×K 조건 없이 "
            "거래대금·시총·MA5 등 나머지 필터만으로 당일 종가 매수합니다. 신호 빈도가 크게 늘 수 있습니다."
        )
    if p.get("mcap_side") == "bottom":
        lines.append(
            f"{len(lines)+1}.  시총 하위 필터 — 매매 전날 기준 cross-section 하위 {p.get('mcap_top_ratio', 0.1)*100:.0f}% "
            "소형주 위주로 진입합니다. 유동성·체결 난이도는 별도 반영되지 않았습니다."
        )
    if p.get("sell_mode") == "same_close":
        lines.append(
            f"{len(lines)+1}.  낮은 MDD — 당일 종가 전량 청산(오버나이트 미보유) 구조라 평가 뱅크롤이 거의 우상향해 "
            "MDD가 낮게 나옵니다. 장중 변동·연속 보유 리스크는 반영되지 않습니다."
        )
        lines.append(
            f"{len(lines)+1}.  체결 가정 — 돌파가 매수·당일 종가 전량 매도, 슬리피지·세금 미반영, 거래비용은 매도 시 "
            f"{_cost_pct(p)} 단순 적용입니다. 복리 배수는 매일 전액 재투자한 이론적 상한이므로 보수적으로 해석해야 합니다."
        )
    else:
        lines.append(
            f"{len(lines)+1}.  낮은 MDD — 매일 아침 전량 현금화(오버나이트 1일 보유) 구조라 평가 뱅크롤이 거의 우상향해 "
            "MDD가 낮게 나옵니다. 장중·연속 보유 리스크는 반영되지 않습니다."
        )
        lines.append(
            f"{len(lines)+1}.  체결 가정 — 익일 시가 전량 체결·슬리피지·세금 미반영, 거래비용은 매도 시 "
            f"{_cost_pct(p)} 단순 적용입니다. "
            "복리 배수는 매일 전액 재투자한 이론적 상한으로, 실제 유동성 제약을 감안하면 보수적으로 해석해야 합니다."
        )
    return "\n\n".join(lines)


def build_report(res: dict, out_name: str, *, title: str, subtitle: str) -> Path:
    print("▶ 차트 생성 중...")
    charts = make_charts(res)
    out = _resolve_out_path(out_name)
    print("▶ PPT 생성 중...")
    build_ppt(
        res, charts, out,
        title=title,
        subtitle=subtitle,
        cond_rows=_build_cond_rows(res["params"]),
        caveats=_build_caveats(res["params"]),
    )
    print(f"✅ 저장 완료: {out}")
    return out


def main() -> None:
    import argparse

    ap = argparse.ArgumentParser()
    ap.add_argument("--mode", choices=["baseline", "snapshot"], default="baseline")
    args = ap.parse_args()

    if args.mode == "snapshot":
        print("▶ [snapshot] 백테스트 실행 중...")
        res = rv.run_backtest(
            start="2024-01-01", end="2025-12-31", data_dir="data/raw",
            output="", seed=42, verbose=False,
        )
        build_report(
            res, "repair_v13_bankroll_backtest_2024_2025.pptx",
            title="repair_v13 백테스트 결과 리포트",
            subtitle="스냅샷 시총 전략 · 뱅크롤 1억 복리 시뮬레이션",
        )
    else:
        import repair_v13_baseline as bl

        print("▶ [baseline] point-in-time 백테스트 실행 중...")
        res = bl.run_baseline(
            start="2020-01-01", end="2026-05-31", data_dir="data/raw",
            seed=42, verbose=False,
        )
        if not res:
            print("❌ 결과 없음")
            return
        build_report(
            res, "repair_v13_baseline_backtest_2020_2026.pptx",
            title="repair_v13 표준(baseline) 백테스트 리포트",
            subtitle="Point-in-time 시총 전략 · 뱅크롤 1억 복리 시뮬레이션",
        )


if __name__ == "__main__":
    main()
